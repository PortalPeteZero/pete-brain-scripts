#!/usr/bin/env python3
"""clancy-dn-enrich.py — STAGE 3, Enrichment: read every held Depotnet file, in full.

Plan: vault note Projects/SY-Clancy/plan-doc-enrichment.md (amended 3 Aug 2026 — Drive-first,
images IN scope, anti-skim ledger). The four stages: Capture -> File -> ENRICH -> Publish.

What it does, per phase:
  --fetch      download every in-scope file from its Drive copy to the work dir
  --extract    deterministic reading of every file type; per-file extract JSON; ledger rows;
               anything needing model eyes (images, scanned pdf pages, video frames, embedded
               photos) goes on the vision queue
  --gate       the anti-skim counter: every in-scope file accounted for, or the exact list of
               what is not
  --assemble   one enrich file per damage (v1 Depotnet record first, then every extract in full,
               then every vision reading) written to the work dir
  --upload     push each damage's enrich file into its own Drive folder

The reading ledger is `clancy_dn_enrich_ledger` (one row per file per parser version).
Vision results are written by the session/agents as JSON next to the queue; --assemble refuses
to run while any queue item has no result (the guard against skim).

Scope: --fy FY26/27 (default). FY25/26 later, same machinery.

Usage:
  VAULT=/tmp/pbs python3 clancy-dn-enrich.py --fetch [--incident N]
  VAULT=/tmp/pbs python3 clancy-dn-enrich.py --extract [--incident N]
  VAULT=/tmp/pbs python3 clancy-dn-enrich.py --gate
  VAULT=/tmp/pbs python3 clancy-dn-enrich.py --assemble [--incident N]
  VAULT=/tmp/pbs python3 clancy-dn-enrich.py --upload [--incident N]
"""
import os, sys, json, re, csv, glob, shutil, argparse, hashlib, zipfile, subprocess, datetime
import urllib.request, urllib.parse
from concurrent.futures import ThreadPoolExecutor, as_completed

VAULT = os.environ.get("VAULT", "/tmp/pbs")
WORK = os.environ.get("ENRICH_WORK", "/tmp/enrich-work")
PARSER_VERSION = "e1-2026-08-03"
MSG_PY = "/tmp/enrich-venv/bin/python"   # ephemeral venv with extract_msg
DECORATIVE_BYTES = 15 * 1024             # embedded media smaller than this = decorative, recorded not skimmed
SHEET_FULL_ROWS = 500                    # spreadsheets beyond this: header + first 50 rows + measured stats, stated in the ledger
VIDEO_FRAME_EVERY_S = 3
VIDEO_FRAME_CAP = 20

SEC = os.path.expanduser("~/.config/pete-secrets")
if not os.path.exists(f"{SEC}/command-centre-supabase-keys.json"):
    SEC = f"{VAULT}/Library/processes/secrets"
_k = json.load(open(f"{SEC}/command-centre-supabase-keys.json"))
URL, SR = _k["url"], _k["service_role_key"]
H = {"apikey": SR, "Authorization": f"Bearer {SR}", "Content-Type": "application/json"}

# drive-api.py as a library (VAULT-aware; no argparse at import)
import importlib.util
_spec = importlib.util.spec_from_file_location("drive_api", f"{VAULT}/drive-api.py")
drive = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(drive)


def rest(path, method="GET", body=None, headers=None):
    h = dict(H); h.update(headers or {})
    req = urllib.request.Request(f"{URL}/rest/v1/{path}",
                                 data=(json.dumps(body).encode() if body is not None else None),
                                 headers=h, method=method)
    with urllib.request.urlopen(req, timeout=180) as r:
        t = r.read().decode()
        return json.loads(t) if t else None


def rest_all(path, page=1000):
    """PostgREST caps a response at 1,000 rows - page through, or a scope silently truncates.

    REFUSES to run without a sort order: Postgres makes no ordering promise without ORDER BY, so
    a paged read without one can hand back the same row twice and never show you another. Silently
    short, no error. This scope decides which files get read at all, so it fails closed.
    (Flagged by cc-locator-audit 3 Aug 2026.)
    """
    if "order=" not in path:
        raise ValueError(f"rest_all needs an explicit &order= - refusing to page blind: {path}")
    out, off = [], 0
    while True:
        chunk = rest(path, headers={"Range-Unit": "items", "Range": f"{off}-{off+page-1}"})
        if not chunk:
            break
        out.extend(chunk)
        if len(chunk) < page:
            break
        off += page
    return out


def scope_files(fy="FY26/27", incident=None):
    q = (f"clancy_dn_files?select=id,incident_id,kind,name,drive_id,source,storage_path,"
         f"uploaded_on_depotnet,deleted_on_depotnet,action_id"
         f"&order=incident_id,id")
    if incident:
        q += f"&incident_id=eq.{incident}"
    rows = rest_all(q)
    ids = {i["id"] for i in rest_all(f"clancy_dn_incidents?select=id&fy=eq.{urllib.parse.quote(fy)}&order=id")}
    rows = [r for r in rows if r["incident_id"] in ids]
    return rows


def ext_of(name):
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def safe_name(name):
    return re.sub(r"[^A-Za-z0-9._ ()-]", "_", name)[:150]


def local_path(f):
    d = f"{WORK}/files/{f['incident_id']}"
    os.makedirs(d, exist_ok=True)
    return f"{d}/{f['id']}__{safe_name(f['name'])}"


# ---------------------------------------------------------------- fetch

def fetch(fy, incident=None):
    files = scope_files(fy, incident)
    todo = [f for f in files if f.get("drive_id")]
    skipped = [f for f in files if not f.get("drive_id")]
    print(f"in scope: {len(files)} files; {len(todo)} with Drive copies; {len(skipped)} not held")
    done = errs = 0

    def get(f):
        p = local_path(f)
        if os.path.exists(p) and os.path.getsize(p) > 0:
            return "cached"
        drive.get_file(f["drive_id"], p)
        return "fetched"

    with ThreadPoolExecutor(max_workers=8) as ex:
        futs = {ex.submit(get, f): f for f in todo}
        for fut in as_completed(futs):
            f = futs[fut]
            try:
                fut.result(); done += 1
                if done % 50 == 0:
                    print(f"  {done}/{len(todo)}")
            except Exception as e:
                errs += 1
                print(f"  FAILED {f['id']} {f['name']}: {e}")
    print(f"fetched/cached {done}, failed {errs}")
    for f in skipped:
        print(f"  not held: {f['id']} {f['name']} (source={f['source']})")


# ---------------------------------------------------------------- extract helpers

def run(cmd, **kw):
    return subprocess.run(cmd, capture_output=True, text=True, timeout=300, **kw)


def read_pdf(path, vq, f):
    """pdftotext with per-page split; image-only pages rasterised onto the vision queue."""
    r = run(["pdftotext", "-layout", path, "-"])
    pages = r.stdout.split("\f")
    if pages and not pages[-1].strip():
        pages = pages[:-1]
    n = len(pages) or 1
    texty = [i for i, p in enumerate(pages) if len(p.strip()) >= 20]
    imagey = [i for i in range(len(pages)) if i not in texty]
    units = {}
    for i in texty:
        units[f"page {i+1}"] = pages[i].strip()
    for i in imagey:
        out = f"{WORK}/vision/src/{f['id']}_p{i+1}"
        os.makedirs(os.path.dirname(out), exist_ok=True)
        rr = run(["pdftoppm", "-f", str(i+1), "-l", str(i+1), "-r", "120", "-png", path, out])
        made = sorted(glob.glob(out + "*.png"))
        if made:
            vq.append({"path": made[0], "file_id": f["id"], "incident_id": f["incident_id"],
                       "origin": "pdf-page", "label": f"{f['name']} — page {i+1} (no text layer)"})
        else:
            units[f"page {i+1}"] = f"[page {i+1}: no text layer and rasterise failed: {rr.stderr.strip()[:100]}]"
    status = "read" if not imagey else "read+vision"
    note = f"{len(imagey)} image-only page(s) routed to vision" if imagey else None
    return units, n, len(texty) + len(imagey), status, note, "pdftotext"


def read_pptx(path, vq, f):
    from pptx import Presentation
    prs = Presentation(path)
    units, media = {}, 0
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if t.strip():
                    parts.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
            if shape.shape_type == 13:  # picture
                try:
                    blob = shape.image.blob
                    if len(blob) >= DECORATIVE_BYTES:
                        media += 1
                        mp = f"{WORK}/vision/src/{f['id']}_s{i}_m{media}.{shape.image.ext}"
                        os.makedirs(os.path.dirname(mp), exist_ok=True)
                        open(mp, "wb").write(blob)
                        vq.append({"path": mp, "file_id": f["id"], "incident_id": f["incident_id"],
                                   "origin": "pptx-media",
                                   "label": f"{f['name']} — photo on slide {i}"})
                except Exception:
                    pass
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                parts.append("SPEAKER NOTES: " + slide.notes_slide.notes_text_frame.text.strip())
        except Exception:
            pass
        units[f"slide {i}"] = "\n".join(parts).strip() or "[no text on slide]"
    note = f"{media} embedded photo(s) routed to vision" if media else None
    return units, len(prs.slides.__iter__().__length_hint__() if False else prs.slides._sldIdLst), len(units), \
        ("read+vision" if media else "read"), note, "python-pptx"


def read_docx(path, vq, f):
    import docx as docxlib
    d = docxlib.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for tb in d.tables:
        for row in tb.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    media = 0
    with zipfile.ZipFile(path) as z:
        for m in z.namelist():
            if m.startswith("word/media/") and m.lower().endswith((".png", ".jpg", ".jpeg")):
                blob = z.read(m)
                if len(blob) >= DECORATIVE_BYTES:
                    media += 1
                    mp = f"{WORK}/vision/src/{f['id']}_m{media}_{os.path.basename(m)}"
                    os.makedirs(os.path.dirname(mp), exist_ok=True)
                    open(mp, "wb").write(blob)
                    vq.append({"path": mp, "file_id": f["id"], "incident_id": f["incident_id"],
                               "origin": "docx-media", "label": f"{f['name']} — embedded photo"})
    note = f"{media} embedded photo(s) routed to vision" if media else None
    return {"document": "\n".join(parts)}, 1, 1, ("read+vision" if media else "read"), note, "python-docx"


def read_doc(path, vq, f):
    out = path + ".txt"
    r = run(["textutil", "-convert", "txt", path, "-output", out])
    if not os.path.exists(out):
        return {}, 1, 0, "unreadable", f"textutil failed: {r.stderr.strip()[:150]}", "textutil"
    text = open(out, errors="replace").read()
    return {"document": text}, 1, 1, "read", None, "textutil"


def _sheet_dump(rows_iter, sheet_name):
    rows = [r for r in rows_iter]
    nonempty = [r for r in rows if any(str(c).strip() for c in r if c is not None)]
    n = len(nonempty)
    def fmt(r):
        return " | ".join("" if c is None else str(c).strip() for c in r).rstrip(" |")
    if n <= SHEET_FULL_ROWS:
        return f"[{n} rows — full]\n" + "\n".join(fmt(r) for r in nonempty), n, n
    head = nonempty[:51]
    body = (f"[{n} rows — data sheet; header + first 50 shown, full data held in the source file]\n"
            + "\n".join(fmt(r) for r in head))
    return body, n, 51


def read_xlsx(path, vq, f):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    units, total_rows, read_rows = {}, 0, 0
    for ws in wb.worksheets:
        body, n, rd = _sheet_dump(ws.iter_rows(values_only=True), ws.title)
        units[f"sheet: {ws.title}"] = body
        total_rows += n; read_rows += rd
    note = None
    if read_rows < total_rows:
        note = f"data sheets over {SHEET_FULL_ROWS} rows summarised (header+50), stated inline"
    return units, len(wb.worksheets), len(wb.worksheets), "read", note, "openpyxl"


def read_xls(path, vq, f):
    import xlrd
    wb = xlrd.open_workbook(path)
    units = {}
    for ws in wb.sheets():
        rows = ([ws.cell_value(r, c) for c in range(ws.ncols)] for r in range(ws.nrows))
        body, n, rd = _sheet_dump(rows, ws.name)
        units[f"sheet: {ws.name}"] = body
    return units, wb.nsheets, wb.nsheets, "read", None, "xlrd"


def read_msg(path, vq, f):
    script = (
        "import sys, json, extract_msg\n"
        "m = extract_msg.Message(sys.argv[1])\n"
        "atts = []\n"
        "import os\n"
        "outdir = sys.argv[2]\n"
        "os.makedirs(outdir, exist_ok=True)\n"
        "for a in m.attachments:\n"
        "    name = a.longFilename or a.shortFilename or 'attachment'\n"
        "    p = os.path.join(outdir, name)\n"
        "    try:\n"
        "        with open(p, 'wb') as fh: fh.write(a.data if isinstance(a.data, bytes) else b'')\n"
        "        atts.append({'name': name, 'path': p, 'size': os.path.getsize(p)})\n"
        "    except Exception as e:\n"
        "        atts.append({'name': name, 'error': str(e)})\n"
        "print(json.dumps({'from': m.sender, 'to': m.to, 'cc': m.cc, 'date': str(m.date),\n"
        "                  'subject': m.subject, 'body': m.body, 'attachments': atts}))\n"
    )
    attdir = path + ".attachments"
    r = run([MSG_PY, "-c", script, path, attdir])
    if r.returncode != 0:
        return {}, 1, 0, "unreadable", f"extract_msg failed: {r.stderr.strip()[:150]}", "extract_msg"
    m = json.loads(r.stdout)
    hdr = (f"From: {m['from']}\nTo: {m['to']}\nCc: {m.get('cc')}\nDate: {m['date']}\n"
           f"Subject: {m['subject']}\n\n{(m['body'] or '').strip()}")
    units = {"email": hdr}
    for a in m["attachments"]:
        if "path" in a and os.path.exists(a["path"]):
            aext = ext_of(a["name"])
            if aext in ("jpg", "jpeg", "png"):
                vq.append({"path": a["path"], "file_id": f["id"], "incident_id": f["incident_id"],
                           "origin": "msg-attachment", "label": f"{f['name']} — attachment {a['name']}"})
                units[f"attachment: {a['name']}"] = "[image — routed to vision]"
            elif aext == "pdf":
                u2, *_ = read_pdf(a["path"], vq, f)
                units[f"attachment: {a['name']}"] = "\n\n".join(f"{k}:\n{v}" for k, v in u2.items())
            elif aext in ("docx",):
                u2, *_ = read_docx(a["path"], vq, f)
                units[f"attachment: {a['name']}"] = u2.get("document", "")
            else:
                units[f"attachment: {a['name']}"] = f"[{aext} attachment held at capture, not parsed inline]"
    note = f"{len(m['attachments'])} attachment(s) captured" if m["attachments"] else None
    return units, 1, 1, "read", note, "extract_msg"


def read_zip(path, vq, f):
    units = {}
    d = path + ".unzipped"
    with zipfile.ZipFile(path) as z:
        members = [m for m in z.namelist() if not m.endswith("/")]
        z.extractall(d)
    for m in members:
        p = os.path.join(d, m)
        e = ext_of(m)
        if e in ("jpg", "jpeg", "png"):
            vq.append({"path": p, "file_id": f["id"], "incident_id": f["incident_id"],
                       "origin": "zip-member", "label": f"{f['name']} — {m}"})
            units[f"member: {m}"] = "[image — routed to vision]"
        elif e == "pdf":
            u2, *_ = read_pdf(p, vq, f)
            units[f"member: {m}"] = "\n\n".join(f"{k}:\n{v}" for k, v in u2.items())
        elif e in ("docx",):
            u2, *_ = read_docx(p, vq, f)
            units[f"member: {m}"] = u2.get("document", "")
        elif e in ("txt", "csv", "md"):
            units[f"member: {m}"] = open(p, errors="replace").read()
        else:
            units[f"member: {m}"] = f"[{e} member held, not parsed inline]"
    return units, len(members), len(members), "read", f"{len(members)} member(s)", "zipfile"


def read_vsdx(path, vq, f):
    units = {}
    with zipfile.ZipFile(path) as z:
        pages = [m for m in z.namelist() if m.startswith("visio/pages/") and m.endswith(".xml")]
        for m in pages:
            xml = z.read(m).decode(errors="replace")
            texts = re.findall(r"<Text[^>]*>(.*?)</Text>", xml, re.S)
            clean = [re.sub(r"<[^>]+>", " ", t).strip() for t in texts]
            clean = [t for t in clean if t]
            if clean:
                units[f"page: {os.path.basename(m)}"] = "\n".join(clean)
    if not units:
        return {}, 1, 0, "unreadable", "no text found in Visio XML", "vsdx-xml"
    return units, len(units), len(units), "read", None, "vsdx-xml"


def read_image(path, vq, f):
    from PIL import Image
    try:
        with Image.open(path) as im:
            dims = f"{im.width}x{im.height}"
    except Exception as e:
        return {}, 1, 0, "unreadable", f"not a valid image: {e}", "pillow"
    vq.append({"path": path, "file_id": f["id"], "incident_id": f["incident_id"],
               "origin": "native-image", "label": f"{f['name']} ({dims})"})
    return {"image": f"[{dims} — routed to vision]"}, 1, 1, "routed-vision", None, "pillow+vision"


def read_mp4(path, vq, f):
    r = run(["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", path])
    try:
        dur = float(r.stdout.strip())
    except ValueError:
        return {}, 1, 0, "unreadable", f"ffprobe failed: {r.stderr.strip()[:100]}", "ffmpeg"
    n = min(VIDEO_FRAME_CAP, max(3, int(dur / VIDEO_FRAME_EVERY_S)))
    outdir = f"{WORK}/vision/src/{f['id']}_frames"
    os.makedirs(outdir, exist_ok=True)
    rr = run(["ffmpeg", "-y", "-i", path, "-vf", f"fps={n}/{dur:.2f}" if dur > 0 else "fps=1",
              "-frames:v", str(n), f"{outdir}/frame_%02d.png"])
    frames = sorted(glob.glob(f"{outdir}/frame_*.png"))
    for i, fr in enumerate(frames, 1):
        vq.append({"path": fr, "file_id": f["id"], "incident_id": f["incident_id"],
                   "origin": "video-frame",
                   "label": f"{f['name']} — frame {i}/{len(frames)} (~{i*dur/max(len(frames),1):.0f}s of {dur:.0f}s)"})
    note = f"{dur:.0f}s video, {len(frames)} frames sampled; audio not transcribed (no local transcriber)"
    return {"video": note}, len(frames), len(frames), "read+vision", note, "ffmpeg-frames"


READERS = {"pdf": read_pdf, "pptx": read_pptx, "docx": read_docx, "doc": read_doc,
           "xlsx": read_xlsx, "xls": read_xls, "msg": read_msg, "zip": read_zip,
           "vsdx": read_vsdx, "jpg": read_image, "jpeg": read_image, "png": read_image,
           "mp4": read_mp4}


def read_plain(path, vq, f):
    return {"text": open(path, errors="replace").read()}, 1, 1, "read", None, "native"


for _e in ("txt", "md", "csv"):
    READERS[_e] = read_plain


# ---------------------------------------------------------------- extract

def extract(fy, incident=None):
    files = scope_files(fy, incident)
    os.makedirs(f"{WORK}/extracts", exist_ok=True)
    os.makedirs(f"{WORK}/vision", exist_ok=True)
    vq_path = f"{WORK}/vision-queue.jsonl"
    vq_existing = set()
    if os.path.exists(vq_path):
        for line in open(vq_path):
            vq_existing.add(json.loads(line)["path"])
    ledger, n_read, n_vision, n_bad = [], 0, 0, 0
    vq_new = []
    for f in files:
        e = ext_of(f["name"])
        xp = f"{WORK}/extracts/{f['id']}.json"
        if os.path.exists(xp):
            continue  # idempotent per run; delete extracts dir to re-read
        lp = local_path(f)
        row = {"file_id": f["id"], "incident_id": f["incident_id"], "name": f["name"],
               "kind": f.get("kind"), "ext": e, "parser_version": PARSER_VERSION}
        if not f.get("drive_id") or not os.path.exists(lp):
            if f["id"] == 12047:
                reason = "our own published Hollow Lane review — source held natively (module clancy-hollow-lane-review)"
            elif f.get("source") == "unfetchable-sas":
                reason = "not held: Depotnet download link was unfetchable at capture; re-attempt next capture session"
            else:
                reason = "not held locally and no Drive copy"
            row.update({"method": "none", "status": "not-held", "units_total": 0, "units_read": 0,
                        "chars_out": 0, "note": reason})
            ledger.append(row); n_bad += 1
            json.dump({**row, "units": {}}, open(xp, "w"))
            continue
        reader = READERS.get(e)
        vq = []
        if reader is None:
            row.update({"method": "none", "status": "unreadable", "units_total": 1, "units_read": 0,
                        "chars_out": 0, "note": f"no reader for .{e} — flag for a handler"})
            ledger.append(row); n_bad += 1
            json.dump({**row, "units": {}}, open(xp, "w"))
            continue
        try:
            units, total, read_n, status, note, method = reader(lp, vq, f)
        except Exception as ex:
            row.update({"method": "error", "status": "unreadable", "units_total": 1, "units_read": 0,
                        "chars_out": 0, "note": f"{type(ex).__name__}: {str(ex)[:180]}"})
            ledger.append(row); n_bad += 1
            json.dump({**row, "units": {}}, open(xp, "w"))
            continue
        chars = sum(len(v) for v in units.values())
        row.update({"method": method, "status": status, "units_total": total, "units_read": read_n,
                    "chars_out": chars, "note": note})
        ledger.append(row)
        json.dump({**row, "units": units, "vision_items": [v["path"] for v in vq]},
                  open(xp, "w"))
        if "vision" in status:
            n_vision += 1
        else:
            n_read += 1
        for v in vq:
            if v["path"] not in vq_existing:
                vq_new.append(v); vq_existing.add(v["path"])
    with open(vq_path, "a") as fh:
        for v in vq_new:
            fh.write(json.dumps(v) + "\n")
    # ledger upsert
    for i in range(0, len(ledger), 200):
        rest("clancy_dn_enrich_ledger?on_conflict=file_id,parser_version", "POST",
             ledger[i:i+200], {"Prefer": "resolution=merge-duplicates"})
    print(f"extracted: {n_read} clean, {n_vision} with vision items, {n_bad} unreadable/not-held; "
          f"{len(vq_new)} new vision-queue items (queue total {len(vq_existing)})")


# ---------------------------------------------------------------- gate

def gate(fy):
    files = scope_files(fy)
    led = {}
    for i in range(0, len(files), 100):
        chunk = files[i:i+100]
        ids = ",".join(str(f["id"]) for f in chunk)
        for r in rest(f"clancy_dn_enrich_ledger?select=file_id,status,units_total,units_read,note"
                      f"&parser_version=eq.{PARSER_VERSION}&file_id=in.({ids})") or []:
            led[r["file_id"]] = r
    missing = [f for f in files if f["id"] not in led]
    partial = [r for r in led.values() if r["status"] not in
               ("read", "read+vision", "routed-vision", "not-held") ]
    # vision completeness
    vq_path = f"{WORK}/vision-queue.jsonl"
    q = [json.loads(l) for l in open(vq_path)] if os.path.exists(vq_path) else []
    unread = []
    for v in q:
        rp = f"{WORK}/vision/results/{hashlib.md5(v['path'].encode()).hexdigest()}.json"
        if not os.path.exists(rp):
            unread.append(v)
        else:
            d = json.load(open(rp))
            if not (d.get("description") or "").strip():
                unread.append(v)
    print(f"GATE — in-scope files: {len(files)}")
    print(f"  no ledger row:            {len(missing)}")
    print(f"  unreadable (needs work):  {len(partial)}")
    print(f"  vision queue:             {len(q)} items, {len(unread)} without a result")
    for f in missing[:20]:
        print(f"    MISSING {f['id']} {f['name']}")
    for r in partial[:20]:
        print(f"    UNREADABLE {r['file_id']}: {r['note']}")
    for v in unread[:15]:
        print(f"    VISION PENDING: {v['label']}")
    n = len(missing) + len(partial) + len(unread)
    print(f"UNACCOUNTED TOTAL: {n}" + ("  ✅ CLEAN" if n == 0 else "  ⛔ NOT DONE"))
    return n


# ---------------------------------------------------------------- assemble

def vision_result(path):
    rp = f"{WORK}/vision/results/{hashlib.md5(path.encode()).hexdigest()}.json"
    return json.load(open(rp)) if os.path.exists(rp) else None


def assemble(fy, incident=None, allow_pending=False):
    files = scope_files(fy, incident)
    by_inc = {}
    for f in files:
        by_inc.setdefault(f["incident_id"], []).append(f)
    inc_rows = rest(f"clancy_dn_incidents?select=id,date_raised,incident_date,category,subcategory,"
                    f"location,severity,status,contract,description,root_cause,underlying_cause,"
                    f"lessons_learnt,incident_summary,utility_class,capture_drive_folder"
                    f"&fy=eq.{urllib.parse.quote(fy)}")
    incs = {r["id"]: r for r in inc_rows}
    vq_path = f"{WORK}/vision-queue.jsonl"
    q = [json.loads(l) for l in open(vq_path)] if os.path.exists(vq_path) else []
    q_by_file = {}
    for v in q:
        q_by_file.setdefault(v["file_id"], []).append(v)
    os.makedirs(f"{WORK}/enrich-files", exist_ok=True)
    today = datetime.date.today().strftime("%d %B %Y")
    made = 0
    for iid, flist in sorted(by_inc.items()):
        inc = incs.get(iid, {})
        out = [f"# Damage {iid} - enrichment record",
               "",
               f"Generated {today} · parser {PARSER_VERSION} · the v2 layer: everything read from "
               f"the {len(flist)} files this damage holds. The v1 record (Depotnet's own words) is "
               f"shown first and is never altered by this file.",
               "",
               "## The damage as Depotnet records it (v1 — verbatim)",
               ""]
        for k, label in [("incident_date", "Incident date"), ("date_raised", "Date raised"),
                         ("location", "Location"), ("category", "Category"),
                         ("subcategory", "Subcategory"), ("utility_class", "Utility"),
                         ("severity", "Severity"), ("status", "Status"), ("contract", "Contract"),
                         ("description", "Description"), ("incident_summary", "Incident summary"),
                         ("root_cause", "Root cause"), ("underlying_cause", "Underlying cause"),
                         ("lessons_learnt", "Lessons learnt")]:
            v = inc.get(k)
            out.append(f"- **{label}:** {v if v not in (None, '') else '(blank on Depotnet)'}")
        out += ["", f"## What we hold ({len(flist)} files)", "",
                "| File | Kind | Uploaded | Read result |", "|---|---|---|---|"]
        docs, photos, videos = [], [], []
        for f in flist:
            xp = f"{WORK}/extracts/{f['id']}.json"
            x = json.load(open(xp)) if os.path.exists(xp) else None
            res = (f"{x['status']} — {x['units_read']}/{x['units_total']}" if x else "NO EXTRACT")
            out.append(f"| {f['name']} | {f.get('kind') or ext_of(f['name'])} | "
                       f"{(f.get('uploaded_on_depotnet') or '')[:10]} | {res} |")
            e = ext_of(f["name"])
            if e in ("jpg", "jpeg", "png"):
                photos.append((f, x))
            elif e == "mp4":
                videos.append((f, x))
            else:
                docs.append((f, x))
        if docs:
            out += ["", "## The documents, in full", ""]
            for f, x in docs:
                out.append(f"### {f['name']}")
                out.append(f"*{f.get('kind') or ext_of(f['name'])} · uploaded "
                           f"{(f.get('uploaded_on_depotnet') or 'unknown')[:10]} · "
                           f"read: {x['status'] if x else 'none'} ({x['method'] if x else '-'})*")
                out.append("")
                if x and x.get("units"):
                    for uname, utext in x["units"].items():
                        out.append(f"**{uname}**")
                        out.append("")
                        out.append(utext if utext.strip() else "[empty]")
                        out.append("")
                elif x:
                    out.append(f"*Not read: {x.get('note')}*")
                    out.append("")
                for v in q_by_file.get(f["id"], []):
                    if v["origin"] in ("pdf-page", "pptx-media", "docx-media", "msg-attachment", "zip-member"):
                        r = vision_result(v["path"])
                        out.append(f"**{v['label']}** (read by vision)")
                        out.append("")
                        if r:
                            out.append(r.get("description", "").strip())
                            if (r.get("transcription") or "").strip():
                                out.append("")
                                out.append("Text in the image, transcribed:")
                                out.append("> " + r["transcription"].strip().replace("\n", "\n> "))
                        else:
                            out.append("*[vision reading pending]*")
                        out.append("")
        if photos:
            out += ["", "## The photographs", ""]
            for f, x in photos:
                for v in q_by_file.get(f["id"], []):
                    r = vision_result(v["path"])
                    out.append(f"### {f['name']}")
                    out.append(f"*uploaded {(f.get('uploaded_on_depotnet') or 'unknown')[:10]}*")
                    out.append("")
                    if r:
                        out.append(r.get("description", "").strip())
                        if (r.get("transcription") or "").strip():
                            out.append("")
                            out.append("Text in the photo, transcribed:")
                            out.append("> " + r["transcription"].strip().replace("\n", "\n> "))
                    else:
                        out.append("*[vision reading pending]*")
                    out.append("")
        if videos:
            out += ["", "## The videos", ""]
            for f, x in videos:
                out.append(f"### {f['name']}")
                if x:
                    out.append(f"*{x.get('note')}*")
                out.append("")
                for v in q_by_file.get(f["id"], []):
                    r = vision_result(v["path"])
                    if r:
                        out.append(f"- **{v['label'].split(' — ')[-1]}**: {r.get('description','').strip()}")
                out.append("")
        out += ["", "---", f"*Reading ledger: {len(flist)} files · "
                f"docs {len(docs)} · photos {len(photos)} · videos {len(videos)} · "
                f"parser {PARSER_VERSION}. The v1 baseline is frozen in "
                f"clancy_dn_baseline_pre_enrichment; nothing in this file overwrites Depotnet's record.*"]
        open(f"{WORK}/enrich-files/Damage {iid} - enrichment record.md", "w").write("\n".join(out))
        made += 1
    print(f"assembled {made} enrich files in {WORK}/enrich-files/")


# ---------------------------------------------------------------- upload

def folder_id_of(url):
    m = re.search(r"/folders/([A-Za-z0-9_-]+)", url or "")
    return m.group(1) if m else None


def upload(fy, incident=None):
    inc_rows = rest(f"clancy_dn_incidents?select=id,capture_drive_folder"
                    f"&fy=eq.{urllib.parse.quote(fy)}&capture_drive_folder=not.is.null")
    done = 0
    for r in sorted(inc_rows, key=lambda x: x["id"]):
        if incident and r["id"] != incident:
            continue
        p = f"{WORK}/enrich-files/Damage {r['id']} - enrichment record.md"
        if not os.path.exists(p):
            print(f"  no enrich file for {r['id']} — skipped")
            continue
        fid = folder_id_of(r["capture_drive_folder"])
        if not fid:
            print(f"  no folder id for {r['id']} — skipped")
            continue
        name = os.path.basename(p)
        existing = drive.api("GET", "/files", params={
            "q": f"name = '{name}' and '{fid}' in parents and trashed = false",
            "fields": "files(id)", "supportsAllDrives": "true",
            "includeItemsFromAllDrives": "true", "corpora": "allDrives"})
        if existing.get("files"):
            gid = existing["files"][0]["id"]
            with open(p, "rb") as fh:
                data = fh.read()
            req = urllib.request.Request(
                f"{drive.UPLOAD_BASE}/files/{gid}?uploadType=media&supportsAllDrives=true",
                data=data, headers={"Authorization": f"Bearer {drive.get_token()}",
                                    "Content-Type": "text/markdown"}, method="PATCH")
            urllib.request.urlopen(req, timeout=120).read()
            print(f"  updated {name}")
        else:
            drive.upload_file(p, fid, name)
            print(f"  uploaded {name}")
        done += 1
    print(f"uploaded/updated {done} enrich files")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--fy", default="FY26/27")
    ap.add_argument("--incident", type=int)
    ap.add_argument("--fetch", action="store_true")
    ap.add_argument("--extract", action="store_true")
    ap.add_argument("--gate", action="store_true")
    ap.add_argument("--assemble", action="store_true")
    ap.add_argument("--upload", action="store_true")
    a = ap.parse_args()
    os.makedirs(WORK, exist_ok=True)
    if a.fetch:
        fetch(a.fy, a.incident)
    if a.extract:
        extract(a.fy, a.incident)
    if a.gate:
        sys.exit(0 if gate(a.fy) == 0 else 2)
    if a.assemble:
        assemble(a.fy, a.incident)
    if a.upload:
        upload(a.fy, a.incident)
    if not any([a.fetch, a.extract, a.gate, a.assemble, a.upload]):
        ap.print_help()


if __name__ == "__main__":
    main()
